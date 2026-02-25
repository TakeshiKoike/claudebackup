# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

pptx_path = r"C:\Users\kokek\OneDrive\デスクトップ\koike  KES2025_Nursing_Metaverse_Presentation.pptx"
out_path = r"C:\Users\kokek\OneDrive\デスクトップ\koike  KES2025_Nursing_Metaverse_Presentation_JP.pptx"

prs = Presentation(pptx_path)

# Translation dictionary: English -> Japanese (run-level text)
translations = {
    # Slide 1
    "Metaverse-Based Clinical Training": "メタバースを活用した臨床トレーニング",
    "in Nursing Education": "看護教育における",
    "Utilization of Computer Graphics and Generative AI": "コンピュータグラフィックスと生成AIの活用",
    "1) Seirei Christopher University": "1) 聖隷クリストファー大学",
    "2) Osaka Metropolitan University Graduate School of Informatics": "2) 大阪公立大学大学院情報学研究科",
    "KES 2025 - 29th International Conference": "KES 2025 - 第29回国際会議",
    "Takeshi Koike ": "小池 武嗣 ",
    " Yukie Majima ": " 真嶋 由貴惠 ",

    # Slide 2
    "Research Background": "研究背景",
    "Global Challenges in Nursing Education": "看護教育におけるグローバルな課題",
    "\u2022 Pandemic constraints on traditional clinical practice": "\u2022 パンデミックによる従来の臨床実習への制約",
    "\u2022 Rapid digital transformation (DX) requirements": "\u2022 急速なデジタルトランスフォーメーション（DX）の要件",
    "\u2022 Need for innovative educational methods": "\u2022 革新的な教育方法の必要性",
    "Current VR Simulation Limitations": "現行VRシミュレーションの限界",
    "\u2022 Insufficient visual realism": "\u2022 視覚的リアリズムの不足",
    "\u2022 Scenario-bound interactions": "\u2022 シナリオに制約されたインタラクション",
    "\u2022 High development costs and time": "\u2022 高い開発コストと時間",

    # Slide 3
    "Challenges in Conventional Nursing Education": "従来の看護教育における課題",
    "Visual Realism": "視覚的リアリズム",
    "Limitations": "の限界",
    "Cannot adequately": "医療現場の複雑さを",
    "reproduce medical": "十分に再現",
    "setting complexity": "できない",
    "Interaction": "インタラクション",
    "Constraints": "の制約",
    "Dependent on": "事前定義された",
    "predefined scenarios,": "シナリオに依存し、",
    "limited flexibility": "柔軟性が限定的",
    "Development": "開発",
    "Efficiency": "効率",
    "Requires extensive": "多大な時間と",
    "time and": "専門知識が",
    "expertise": "必要",
    "Lack of Integrated": "統合的学習の",
    "Learning": "欠如",
    "Fragmented learning": "断片的な学習",
    "approach": "アプローチ",
    "Educational Impacts": "教育への影響",
    "Limited Learning Effectiveness": "限定的な学習効果",
    "Inadequate Practice Preparation": "不十分な実習準備",
    "Widening Educational Disparities": "拡大する教育格差",
    "Reduced Patient Safety": "低下する患者安全",

    # Slide 4
    "Research Purpose & Innovation": "研究目的とイノベーション",
    "Main Objectives": "主な目的",
    "\u2022 Clarify efficient construction methods for nursing education metaverse": "\u2022 看護教育メタバースの効率的な構築方法の明確化",
    "\u2022 Develop effective integration methods with generative AI": "\u2022 生成AIとの効果的な統合方法の開発",
    "\u2022 Create practical implementation guidelines": "\u2022 実践的な導入ガイドラインの作成",
    "Key Innovations": "主要なイノベーション",
    "\u2022 UE5-specific approach for nursing education": "\u2022 看護教育に特化したUE5アプローチ",
    "\u2022 Efficient production workflows for educational institutions": "\u2022 教育機関向けの効率的な制作ワークフロー",
    "\u2022 Deep AI integration with psychological/physiological states": "\u2022 心理的・生理的状態とのAI深層統合",

    # Slide 5
    "Innovative Solution: UE5 + Generative AI": "革新的ソリューション：UE5 + 生成AI",
    "Unreal Engine 5": "Unreal Engine 5",
    "High-Quality 3D Rendering": "高品質3Dレンダリング",
    "\u2022 Detailed medical equipment": "\u2022 精密な医療機器",
    "\u2022 Hospital environments": "\u2022 病院環境",
    "MetaHuman Creator": "MetaHuman Creator",
    "\u2022 Diverse patient representation": "\u2022 多様な患者表現",
    "\u2022 Realistic facial expressions": "\u2022 リアルな表情",
    "Real-time Systems": "リアルタイムシステム",
    "\u2022 Dynamic lighting": "\u2022 ダイナミックライティング",
    "\u2022 Fluid simulations": "\u2022 流体シミュレーション",
    "Generative AI": "生成AI",
    "Dynamic Patient System": "ダイナミック患者システム",
    "\u2022 Natural language responses": "\u2022 自然言語応答",
    "\u2022 Unexpected behaviors": "\u2022 予期しない反応",
    "Individualized Learning": "個別化学習",
    "\u2022 Adaptive responses": "\u2022 適応型応答",
    "\u2022 Performance analytics": "\u2022 パフォーマンス分析",
    "Symptom Progression": "症状の進行",
    "\u2022 Treatment-based changes": "\u2022 治療に基づく変化",
    "\u2022 Real-time feedback": "\u2022 リアルタイムフィードバック",

    # Slide 8
    "System Architecture: 4-Layer Structure": "システムアーキテクチャ：4層構造",
    "Layer 4: Learning Analytics": "レイヤー4：学習分析",
    "Behavior recording | AI feedback | Goal visualization": "行動記録 | AIフィードバック | 目標の可視化",
    "Layer 3: Patient State Management": "レイヤー3：患者状態管理",
    "Physiological parameters | Psychological states | Symptom progression": "生理学的パラメータ | 心理状態 | 症状の進行",
    "Layer 2: AI Integration Middleware": "レイヤー2：AI統合ミドルウェア",
    "REST API (GPT-4, Claude) | JSON data exchange | Response queuing": "REST API (GPT-4, Claude) | JSONデータ交換 | 応答キューイング",
    "Layer 1: UE5 Core Environment": "レイヤー1：UE5コア環境",
    "3D hospital | MetaHuman models | Medical equipment": "3D病院 | MetaHumanモデル | 医療機器",

    # Slide 12
    "Create an image or video with generative AI": "生成AIで画像や動画を作成",
    "Generative AI Technology": "生成AI技術",
    "Create conversations (sentences) with generative AI": "生成AIで会話（文章）を作成",
    "Communicate": "コミュニケーション",
    "Programming the Patient Role in Generative AI": "生成AIにおける患者役のプログラミング",
    "Start creating a digital simulated patient": "デジタル模擬患者の作成を開始",

    # Slide 16 / 21
    "With patients and users": "患者およびユーザーとの",
    "communication": "コミュニケーション",
    "Physical assessment of the patient": "患者のフィジカルアセスメント",
    "Nursing process": "看護過程",
    "Clinical experience": "臨床体験",
    "\u2460 Artificial intelligence AI": "\u2460 人工知能AI",
    "\u2461 Deep fake": "\u2461 ディープフェイク",
    "\u2462 DX simulated patient using voice recognition etc.": "\u2462 音声認識等を活用したDX模擬患者",
    "\u2460 Actual image": "\u2460 実写映像",
    "\u2461 Computer graphics": "\u2461 コンピュータグラフィックス",
    "\u2462 Game fiction": "\u2462 ゲームフィクション",
    "\u2463 DX bedside utilizing live broadcasting etc.": "\u2463 ライブ配信等を活用したDXベッドサイド",
    "\u2462 360 degree shooting": "\u2462 360度撮影",
    "\u2463 DX clinical space utilizing live broadcasting etc.": "\u2463 ライブ配信等を活用したDX臨床空間",
    "\u2460 Simulated electronic medical record": "\u2460 模擬電子カルテ",
    " Visualization of case examples by CG": " CGによる症例の可視化",
    "\u2462 DX patient setting utilizing game fiction, etc.": "\u2462 ゲームフィクション等を活用したDX患者設定",
    "DX in the development of digital simulated patients": "デジタル模擬患者開発におけるDX",
    "PC screen": "PC画面",
    "Multi-display": "マルチディスプレイ",
    "Virtual reality device": "バーチャルリアリティデバイス",
    "Personal computer": "パソコン",
    "School computer": "学校のコンピュータ",
    "Small display": "小型ディスプレイ",
    "Large display": "大型ディスプレイ",
    "Screen projection by projector": "プロジェクターによるスクリーン投影",
    "AR/XR equipment": "AR/XR機器",
    "VR equipment": "VR機器",
    "Projection Mapping": "プロジェクションマッピング",
    "Reciprocity and reality": "双方向性と現実感",
    "Equipment related to virtual environment": "バーチャル環境に関する機器",
    "Virtual environment output device type": "バーチャル環境出力デバイスの種類",
    "Training content and types of virtual environment": "トレーニング内容とバーチャル環境の種類",
    "mixed reality": "複合現実",
    "Metaverse-Based Clinical Training in Nursing Education: ": "看護教育におけるメタバースを活用した臨床トレーニング：",
    "Utilization of Computer Graphics and Generative": "コンピュータグラフィックスと生成AIの",

    # Slide 17 / 22
    "Maternal Nursing": "母性看護",
    "Adult Nursing": "成人看護",
    "Home Nursing": "在宅看護",
    "Incubator for premature babies": "未熟児用保育器",
    "Home-visit nursing": "訪問看護",
    "sickroom": "病室",

    # Slide 27
    "PowerPoint 3-screen video production": "PowerPoint 3画面映像制作",
    "Since it is a projector, there are some issues with brightness and resolution.": "プロジェクターのため、明るさと解像度に課題あり",
    "In large LCD display\nStart researching VR construction!": "大型液晶ディスプレイで\nVR構築の研究を開始！",
    "Built a system that allows several people to experience VR at once with 3 prime projectors + 3 rear projection screens.": "3台のプロジェクター＋3面リアプロジェクションスクリーンで複数人が同時にVR体験できるシステムを構築",
    "Three screens output from a three-screen spanned computer to three projectors": "3画面スパンPCから3台のプロジェクターへ3画面出力",
    "Projection from 3 projectors to 3 rear projection large screens": "3台のプロジェクターから3面リアプロジェクション大型スクリーンへ投影",
    "CAVE-type VR construction completed": "CAVE型VR構築完了",
    "Consideration of large displays": "大型ディスプレイの検討",
    "It turned out that it was still difficult to handle because it was heavy and took up a lot of space.": "重量があり場所を取るため、取り扱いが困難であることが判明",
    "Purchased 3 LG 75-inch 4K displays for about 400,000 yen (including stand)": "LG 75インチ4Kディスプレイ3台を約40万円（スタンド込み）で購入",
    "Brighter and clearer than a projector": "プロジェクターより明るく鮮明",

    # Slide 31
    "Phased Learning Program": "段階的学習プログラム",
    "Basic Level": "基礎レベル",
    "Observation & Communication": "観察とコミュニケーション",
    "\u2022 Basic conversation": "\u2022 基本的な会話",
    "\u2022 Vital signs": "\u2022 バイタルサイン",
    "\u2022 Symptom observation": "\u2022 症状観察",
    "Intermediate Level": "中級レベル",
    "Judgment & Practice": "判断と実践",
    "\u2022 Care planning": "\u2022 ケア計画",
    "\u2022 Equipment operation": "\u2022 機器操作",
    "\u2022 Emergency response": "\u2022 緊急対応",
    "Advanced Level": "上級レベル",
    "Integration & Application": "統合と応用",
    "\u2022 Multiple patients": "\u2022 複数患者",
    "\u2022 Team collaboration": "\u2022 チーム連携",
    "\u2022 Ethical dilemmas": "\u2022 倫理的ジレンマ",
    "AI-Powered Individualized Learning": "AI駆動の個別化学習",
    "Weakness identification | Scenario adjustment | Performance tracking": "弱点の特定 | シナリオ調整 | パフォーマンス追跡",

    # Slide 32
    "Implementation Challenges & Solutions": "実装上の課題と解決策",
    "Technical Challenges": "技術的課題",
    "\u2022 AI hallucination phenomena": "\u2022 AIハルシネーション現象",
    "\u2022 System integration complexity": "\u2022 システム統合の複雑さ",
    "\u2022 Data privacy protection": "\u2022 データプライバシー保護",
    "\u2022 Hardware requirements": "\u2022 ハードウェア要件",
    "Technical Solutions": "技術的解決策",
    "\u2022 Multi-layer verification systems": "\u2022 多層検証システム",
    "\u2022 Phased implementation approach": "\u2022 段階的導入アプローチ",
    "\u2022 Privacy protection frameworks": "\u2022 プライバシー保護フレームワーク",
    "\u2022 Cloud-based infrastructure": "\u2022 クラウドベースのインフラ",
    "Organizational Challenges": "組織的課題",
    "\u2022 Faculty training needs": "\u2022 教員研修の必要性",
    "\u2022 Initial investment costs": "\u2022 初期投資コスト",
    "\u2022 Change management": "\u2022 変更管理",
    "Organizational Solutions": "組織的解決策",
    "\u2022 Phased training programs": "\u2022 段階的研修プログラム",
    "\u2022 Strategic investment planning": "\u2022 戦略的投資計画",
    "\u2022 Stakeholder engagement": "\u2022 ステークホルダーの関与",

    # Slide 37
    "1. Far away from the palace, far away from the guidance of the introduction": "1. 遠隔実習指導の導入",
    "Hospital A and University A Relay": "A病院とA大学の中継",
    "University A Hospital": "A大学病院",
    "Send and receive credit PC": "送受信用PC",
    "LTE or ": "LTEまたは",
    "Photography camera": "撮影用カメラ",
    "A University": "A大学",
    "Relay of clinical scenes": "臨床場面の中継",
    "Distance Instruction for Students": "学生への遠隔指導",
    "Teleconferencing": "テレビ会議",
    "2. Introduction of cutting-edge technology into medical education": "2. 医療教育への先端技術の導入",
    "3. Utilization of XR for practical training guidance": "3. 実習指導へのXR活用",
    "4. Application to medical education using existing state-of-the-art technology": "4. 既存の最先端技術を用いた医療教育への応用",
    "Leveraging the Metaverse": "メタバースの活用",
    "Utilization of CG": "CGの活用",
    "Tablet devices": "タブレット端末",
    "Large screen": "大画面",
    "projector": "プロジェクター",
    "Large-format display": "大型ディスプレイ",
    "A: Within the University\nReceiving video as a group": "A：学内\nグループでの映像受信",
    "Metaverse\nplatform": "メタバース\nプラットフォーム",
    "CG software": "CGソフトウェア",
    "High-performance personal computers": "高性能パソコン",
    "On-campus Wi-Fi environment": "学内Wi-Fi環境",
    "Joint development of digital teaching materials": "デジタル教材の共同開発",
    "VR experience at the University of F": "F大学でのVR体験",
    "Integrated at the University of C": "C大学での統合",
    "Leveraging XR\nHolding joint conferences, etc.": "XRの活用\n合同カンファレンスの開催等",
    "Provision of practical training content to desired universities": "希望する大学への実習コンテンツの提供",
    "VR Experience at E University": "E大学でのVR体験",
    "MR experience at the University of D": "D大学でのMR体験",
    "A: Within the University\nIndividual video reception": "A：学内\n個人での映像受信",
    "VR devices ": "VRデバイス ",
    "VR devices": "VRデバイス",
    "Utilization of HTML programs": "HTMLプログラムの活用",
    "Leverage the capture function": "キャプチャ機能の活用",
    "Use of 360-degree cameras": "360度カメラの活用",
    "CG production at University A": "A大学でのCG制作",
    "XR production at University B": "B大学でのXR制作",

    # Slide 38
    "Conclusions": "結論",
    "Key Contributions": "主要な貢献",
    "\u2022 UE5 enables unprecedented clinical environment reproduction": "\u2022 UE5による前例のない臨床環境の再現",
    "\u2022 Generative AI creates dynamic, individualized learning experiences": "\u2022 生成AIによるダイナミックで個別化された学習体験の創出",
    "\u2022 Efficient workflows make implementation feasible for institutions": "\u2022 効率的なワークフローにより教育機関での導入が実現可能に",
    "Future Impact": "将来への影響",
    "\u2022 Fundamental transformation of nursing education paradigm": "\u2022 看護教育パラダイムの根本的な変革",
    "\u2022 Overcome constraints: safety, individualization, repeatability": "\u2022 制約の克服：安全性、個別化、反復可能性",
    "\u2022 Contribute to qualitative improvement and widespread adoption": "\u2022 質的向上と普及への貢献",
    "Nursing education metaverse: Not just technology, but educational transformation": "看護教育メタバース：技術だけでなく、教育の変革",

    # Slide 39
    "Thank You": "ご清聴ありがとうございました",
    "Questions & Discussion": "質疑応答",
    "Contact Information": "連絡先",
    "Takeshi Koike": "小池 武嗣",
    "Seirei Christopher University": "聖隷クリストファー大学",
    "Yukie Majima": "真嶋 由貴惠",
    "Osaka Metropolitan University": "大阪公立大学",

    # Product names that stay as-is (but with Japanese labels in context)
    # iClone, vSim, WiFi, CG, UNITY, TWINMOTION etc. are kept as-is
    "\u00ae for Nursing": "\u00ae for Nursing",
}

def translate_run(run):
    text = run.text
    if text in translations:
        run.text = translations[text]
        return True
    return False

count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if translate_run(run):
                        count += 1
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if translate_run(run):
                                count += 1

prs.save(out_path)
print(f"Translated {count} text runs")
print(f"Saved to: {out_path}")

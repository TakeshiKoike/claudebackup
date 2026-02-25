# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation

pptx_path = r"C:\Users\kokek\OneDrive\デスクトップ\koike  KES2025_Nursing_Metaverse_Presentation.pptx"
out_path = r"C:\Users\kokek\OneDrive\デスクトップ\koike  KES2025_Nursing_Metaverse_Presentation_JP_v2.pptx"

prs = Presentation(pptx_path)

# Translation dictionary: English -> Japanese (run-level text)
# ALL English including product names translated to Japanese
translations = {
    # Slide 1
    "Metaverse-Based Clinical Training": "\u30e1\u30bf\u30d0\u30fc\u30b9\u3092\u6d3b\u7528\u3057\u305f\u81e8\u5e8a\u30c8\u30ec\u30fc\u30cb\u30f3\u30b0",
    "in Nursing Education": "\u770b\u8b77\u6559\u80b2\u306b\u304a\u3051\u308b",
    "Utilization of Computer Graphics and Generative AI": "\u30b3\u30f3\u30d4\u30e5\u30fc\u30bf\u30b0\u30e9\u30d5\u30a3\u30c3\u30af\u30b9\u3068\u751f\u6210AI\u306e\u6d3b\u7528",
    "1) Seirei Christopher University": "1) \u8056\u96b7\u30af\u30ea\u30b9\u30c8\u30d5\u30a1\u30fc\u5927\u5b66",
    "2) Osaka Metropolitan University Graduate School of Informatics": "2) \u5927\u962a\u516c\u7acb\u5927\u5b66\u5927\u5b66\u9662\u60c5\u5831\u5b66\u7814\u7a76\u79d1",
    "KES 2025 - 29th International Conference": "KES 2025 - \u7b2c29\u56de\u56fd\u969b\u4f1a\u8b70",
    "Takeshi Koike ": "\u5c0f\u6c60 \u6b66\u55e3 ",
    " Yukie Majima ": " \u771f\u5d8b \u7531\u8cb4\u60e0 ",

    # Slide 2
    "Research Background": "\u7814\u7a76\u80cc\u666f",
    "Global Challenges in Nursing Education": "\u770b\u8b77\u6559\u80b2\u306b\u304a\u3051\u308b\u30b0\u30ed\u30fc\u30d0\u30eb\u306a\u8ab2\u984c",
    "\u2022 Pandemic constraints on traditional clinical practice": "\u2022 \u30d1\u30f3\u30c7\u30df\u30c3\u30af\u306b\u3088\u308b\u5f93\u6765\u306e\u81e8\u5e8a\u5b9f\u7fd2\u3078\u306e\u5236\u7d04",
    "\u2022 Rapid digital transformation (DX) requirements": "\u2022 \u6025\u901f\u306a\u30c7\u30b8\u30bf\u30eb\u30c8\u30e9\u30f3\u30b9\u30d5\u30a9\u30fc\u30e1\u30fc\u30b7\u30e7\u30f3\uff08DX\uff09\u306e\u8981\u4ef6",
    "\u2022 Need for innovative educational methods": "\u2022 \u9769\u65b0\u7684\u306a\u6559\u80b2\u65b9\u6cd5\u306e\u5fc5\u8981\u6027",
    "Current VR Simulation Limitations": "\u73fe\u884cVR\u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3\u306e\u9650\u754c",
    "\u2022 Insufficient visual realism": "\u2022 \u8996\u899a\u7684\u30ea\u30a2\u30ea\u30ba\u30e0\u306e\u4e0d\u8db3",
    "\u2022 Scenario-bound interactions": "\u2022 \u30b7\u30ca\u30ea\u30aa\u306b\u5236\u7d04\u3055\u308c\u305f\u30a4\u30f3\u30bf\u30e9\u30af\u30b7\u30e7\u30f3",
    "\u2022 High development costs and time": "\u2022 \u9ad8\u3044\u958b\u767a\u30b3\u30b9\u30c8\u3068\u6642\u9593",

    # Slide 3
    "Challenges in Conventional Nursing Education": "\u5f93\u6765\u306e\u770b\u8b77\u6559\u80b2\u306b\u304a\u3051\u308b\u8ab2\u984c",
    "Visual Realism": "\u8996\u899a\u7684\u30ea\u30a2\u30ea\u30ba\u30e0",
    "Limitations": "\u306e\u9650\u754c",
    "Cannot adequately": "\u533b\u7642\u73fe\u5834\u306e\u8907\u96d1\u3055\u3092",
    "reproduce medical": "\u5341\u5206\u306b\u518d\u73fe",
    "setting complexity": "\u3067\u304d\u306a\u3044",
    "Interaction": "\u30a4\u30f3\u30bf\u30e9\u30af\u30b7\u30e7\u30f3",
    "Constraints": "\u306e\u5236\u7d04",
    "Dependent on": "\u4e8b\u524d\u5b9a\u7fa9\u3055\u308c\u305f",
    "predefined scenarios,": "\u30b7\u30ca\u30ea\u30aa\u306b\u4f9d\u5b58\u3057\u3001",
    "limited flexibility": "\u67d4\u8edf\u6027\u304c\u9650\u5b9a\u7684",
    "Development": "\u958b\u767a",
    "Efficiency": "\u52b9\u7387",
    "Requires extensive": "\u591a\u5927\u306a\u6642\u9593\u3068",
    "time and": "\u5c02\u9580\u77e5\u8b58\u304c",
    "expertise": "\u5fc5\u8981",
    "Lack of Integrated": "\u7d71\u5408\u7684\u5b66\u7fd2\u306e",
    "Learning": "\u6b20\u5982",
    "Fragmented learning": "\u65ad\u7247\u7684\u306a\u5b66\u7fd2",
    "approach": "\u30a2\u30d7\u30ed\u30fc\u30c1",
    "Educational Impacts": "\u6559\u80b2\u3078\u306e\u5f71\u97ff",
    "Limited Learning Effectiveness": "\u9650\u5b9a\u7684\u306a\u5b66\u7fd2\u52b9\u679c",
    "Inadequate Practice Preparation": "\u4e0d\u5341\u5206\u306a\u5b9f\u7fd2\u6e96\u5099",
    "Widening Educational Disparities": "\u62e1\u5927\u3059\u308b\u6559\u80b2\u683c\u5dee",
    "Reduced Patient Safety": "\u4f4e\u4e0b\u3059\u308b\u60a3\u8005\u5b89\u5168",

    # Slide 4
    "Research Purpose & Innovation": "\u7814\u7a76\u76ee\u7684\u3068\u30a4\u30ce\u30d9\u30fc\u30b7\u30e7\u30f3",
    "Main Objectives": "\u4e3b\u306a\u76ee\u7684",
    "\u2022 Clarify efficient construction methods for nursing education metaverse": "\u2022 \u770b\u8b77\u6559\u80b2\u30e1\u30bf\u30d0\u30fc\u30b9\u306e\u52b9\u7387\u7684\u306a\u69cb\u7bc9\u65b9\u6cd5\u306e\u660e\u78ba\u5316",
    "\u2022 Develop effective integration methods with generative AI": "\u2022 \u751f\u6210AI\u3068\u306e\u52b9\u679c\u7684\u306a\u7d71\u5408\u65b9\u6cd5\u306e\u958b\u767a",
    "\u2022 Create practical implementation guidelines": "\u2022 \u5b9f\u8df5\u7684\u306a\u5c0e\u5165\u30ac\u30a4\u30c9\u30e9\u30a4\u30f3\u306e\u4f5c\u6210",
    "Key Innovations": "\u4e3b\u8981\u306a\u30a4\u30ce\u30d9\u30fc\u30b7\u30e7\u30f3",
    "\u2022 UE5-specific approach for nursing education": "\u2022 \u770b\u8b77\u6559\u80b2\u306b\u7279\u5316\u3057\u305fUE5\u30a2\u30d7\u30ed\u30fc\u30c1",
    "\u2022 Efficient production workflows for educational institutions": "\u2022 \u6559\u80b2\u6a5f\u95a2\u5411\u3051\u306e\u52b9\u7387\u7684\u306a\u5236\u4f5c\u30ef\u30fc\u30af\u30d5\u30ed\u30fc",
    "\u2022 Deep AI integration with psychological/physiological states": "\u2022 \u5fc3\u7406\u7684\u30fb\u751f\u7406\u7684\u72b6\u614b\u3068\u306eAI\u6df1\u5c64\u7d71\u5408",

    # Slide 5
    "Innovative Solution: UE5 + Generative AI": "\u9769\u65b0\u7684\u30bd\u30ea\u30e5\u30fc\u30b7\u30e7\u30f3\uff1aUE5 + \u751f\u6210AI",
    "Unreal Engine 5": "\u30a2\u30f3\u30ea\u30a2\u30eb\u30a8\u30f3\u30b8\u30f35",
    "High-Quality 3D Rendering": "\u9ad8\u54c1\u8cea3D\u30ec\u30f3\u30c0\u30ea\u30f3\u30b0",
    "\u2022 Detailed medical equipment": "\u2022 \u7cbe\u5bc6\u306a\u533b\u7642\u6a5f\u5668",
    "\u2022 Hospital environments": "\u2022 \u75c5\u9662\u74b0\u5883",
    "MetaHuman Creator": "\u30e1\u30bf\u30d2\u30e5\u30fc\u30de\u30f3 \u30af\u30ea\u30a8\u30a4\u30bf\u30fc",
    "\u2022 Diverse patient representation": "\u2022 \u591a\u69d8\u306a\u60a3\u8005\u8868\u73fe",
    "\u2022 Realistic facial expressions": "\u2022 \u30ea\u30a2\u30eb\u306a\u8868\u60c5",
    "Real-time Systems": "\u30ea\u30a2\u30eb\u30bf\u30a4\u30e0\u30b7\u30b9\u30c6\u30e0",
    "\u2022 Dynamic lighting": "\u2022 \u30c0\u30a4\u30ca\u30df\u30c3\u30af\u30e9\u30a4\u30c6\u30a3\u30f3\u30b0",
    "\u2022 Fluid simulations": "\u2022 \u6d41\u4f53\u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3",
    "Generative AI": "\u751f\u6210AI",
    "Dynamic Patient System": "\u30c0\u30a4\u30ca\u30df\u30c3\u30af\u60a3\u8005\u30b7\u30b9\u30c6\u30e0",
    "\u2022 Natural language responses": "\u2022 \u81ea\u7136\u8a00\u8a9e\u5fdc\u7b54",
    "\u2022 Unexpected behaviors": "\u2022 \u4e88\u671f\u3057\u306a\u3044\u53cd\u5fdc",
    "Individualized Learning": "\u500b\u5225\u5316\u5b66\u7fd2",
    "\u2022 Adaptive responses": "\u2022 \u9069\u5fdc\u578b\u5fdc\u7b54",
    "\u2022 Performance analytics": "\u2022 \u30d1\u30d5\u30a9\u30fc\u30de\u30f3\u30b9\u5206\u6790",
    "Symptom Progression": "\u75c7\u72b6\u306e\u9032\u884c",
    "\u2022 Treatment-based changes": "\u2022 \u6cbb\u7642\u306b\u57fa\u3065\u304f\u5909\u5316",
    "\u2022 Real-time feedback": "\u2022 \u30ea\u30a2\u30eb\u30bf\u30a4\u30e0\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af",

    # Slide 8
    "System Architecture: 4-Layer Structure": "\u30b7\u30b9\u30c6\u30e0\u30a2\u30fc\u30ad\u30c6\u30af\u30c1\u30e3\uff1a4\u5c64\u69cb\u9020",
    "Layer 4: Learning Analytics": "\u30ec\u30a4\u30e4\u30fc4\uff1a\u5b66\u7fd2\u5206\u6790",
    "Behavior recording | AI feedback | Goal visualization": "\u884c\u52d5\u8a18\u9332 | AI\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af | \u76ee\u6a19\u306e\u53ef\u8996\u5316",
    "Layer 3: Patient State Management": "\u30ec\u30a4\u30e4\u30fc3\uff1a\u60a3\u8005\u72b6\u614b\u7ba1\u7406",
    "Physiological parameters | Psychological states | Symptom progression": "\u751f\u7406\u5b66\u7684\u30d1\u30e9\u30e1\u30fc\u30bf | \u5fc3\u7406\u72b6\u614b | \u75c7\u72b6\u306e\u9032\u884c",
    "Layer 2: AI Integration Middleware": "\u30ec\u30a4\u30e4\u30fc2\uff1aAI\u7d71\u5408\u30df\u30c9\u30eb\u30a6\u30a7\u30a2",
    "REST API (GPT-4, Claude) | JSON data exchange | Response queuing": "REST API (GPT-4, Claude) | JSON\u30c7\u30fc\u30bf\u4ea4\u63db | \u5fdc\u7b54\u30ad\u30e5\u30fc\u30a4\u30f3\u30b0",
    "Layer 1: UE5 Core Environment": "\u30ec\u30a4\u30e4\u30fc1\uff1aUE5\u30b3\u30a2\u74b0\u5883",
    "3D hospital | MetaHuman models | Medical equipment": "3D\u75c5\u9662 | \u30e1\u30bf\u30d2\u30e5\u30fc\u30de\u30f3\u30e2\u30c7\u30eb | \u533b\u7642\u6a5f\u5668",

    # Slide 12
    "Create an image or video with generative AI": "\u751f\u6210AI\u3067\u753b\u50cf\u3084\u52d5\u753b\u3092\u4f5c\u6210",
    "Generative AI Technology": "\u751f\u6210AI\u6280\u8853",
    "Create conversations (sentences) with generative AI": "\u751f\u6210AI\u3067\u4f1a\u8a71\uff08\u6587\u7ae0\uff09\u3092\u4f5c\u6210",
    "Communicate": "\u30b3\u30df\u30e5\u30cb\u30b1\u30fc\u30b7\u30e7\u30f3",
    "Programming the Patient Role in Generative AI": "\u751f\u6210AI\u306b\u304a\u3051\u308b\u60a3\u8005\u5f79\u306e\u30d7\u30ed\u30b0\u30e9\u30df\u30f3\u30b0",
    "Start creating a digital simulated patient": "\u30c7\u30b8\u30bf\u30eb\u6a21\u64ec\u60a3\u8005\u306e\u4f5c\u6210\u3092\u958b\u59cb",

    # Slide 16 / 21
    "With patients and users": "\u60a3\u8005\u304a\u3088\u3073\u30e6\u30fc\u30b6\u30fc\u3068\u306e",
    "communication": "\u30b3\u30df\u30e5\u30cb\u30b1\u30fc\u30b7\u30e7\u30f3",
    "Physical assessment of the patient": "\u60a3\u8005\u306e\u30d5\u30a3\u30b8\u30ab\u30eb\u30a2\u30bb\u30b9\u30e1\u30f3\u30c8",
    "Nursing process": "\u770b\u8b77\u904e\u7a0b",
    "Clinical experience": "\u81e8\u5e8a\u4f53\u9a13",
    "\u2460 Artificial intelligence AI": "\u2460 \u4eba\u5de5\u77e5\u80fdAI",
    "\u2461 Deep fake": "\u2461 \u30c7\u30a3\u30fc\u30d7\u30d5\u30a7\u30a4\u30af",
    "\u2462 DX simulated patient using voice recognition etc.": "\u2462 \u97f3\u58f0\u8a8d\u8b58\u7b49\u3092\u6d3b\u7528\u3057\u305fDX\u6a21\u64ec\u60a3\u8005",
    "\u2460 Actual image": "\u2460 \u5b9f\u5199\u6620\u50cf",
    "\u2461 Computer graphics": "\u2461 \u30b3\u30f3\u30d4\u30e5\u30fc\u30bf\u30b0\u30e9\u30d5\u30a3\u30c3\u30af\u30b9",
    "\u2462 Game fiction": "\u2462 \u30b2\u30fc\u30e0\u30d5\u30a3\u30af\u30b7\u30e7\u30f3",
    "\u2463 DX bedside utilizing live broadcasting etc.": "\u2463 \u30e9\u30a4\u30d6\u914d\u4fe1\u7b49\u3092\u6d3b\u7528\u3057\u305fDX\u30d9\u30c3\u30c9\u30b5\u30a4\u30c9",
    "\u2462 360 degree shooting": "\u2462 360\u5ea6\u64ae\u5f71",
    "\u2463 DX clinical space utilizing live broadcasting etc.": "\u2463 \u30e9\u30a4\u30d6\u914d\u4fe1\u7b49\u3092\u6d3b\u7528\u3057\u305fDX\u81e8\u5e8a\u7a7a\u9593",
    "\u2460 Simulated electronic medical record": "\u2460 \u6a21\u64ec\u96fb\u5b50\u30ab\u30eb\u30c6",
    " Visualization of case examples by CG": " CG\u306b\u3088\u308b\u75c7\u4f8b\u306e\u53ef\u8996\u5316",
    "\u2462 DX patient setting utilizing game fiction, etc.": "\u2462 \u30b2\u30fc\u30e0\u30d5\u30a3\u30af\u30b7\u30e7\u30f3\u7b49\u3092\u6d3b\u7528\u3057\u305fDX\u60a3\u8005\u8a2d\u5b9a",
    "DX in the development of digital simulated patients": "\u30c7\u30b8\u30bf\u30eb\u6a21\u64ec\u60a3\u8005\u958b\u767a\u306b\u304a\u3051\u308bDX",
    "PC screen": "PC\u753b\u9762",
    "Multi-display": "\u30de\u30eb\u30c1\u30c7\u30a3\u30b9\u30d7\u30ec\u30a4",
    "Virtual reality device": "\u30d0\u30fc\u30c1\u30e3\u30eb\u30ea\u30a2\u30ea\u30c6\u30a3\u30c7\u30d0\u30a4\u30b9",
    "Personal computer": "\u30d1\u30bd\u30b3\u30f3",
    "School computer": "\u5b66\u6821\u306e\u30b3\u30f3\u30d4\u30e5\u30fc\u30bf",
    "Small display": "\u5c0f\u578b\u30c7\u30a3\u30b9\u30d7\u30ec\u30a4",
    "Large display": "\u5927\u578b\u30c7\u30a3\u30b9\u30d7\u30ec\u30a4",
    "Screen projection by projector": "\u30d7\u30ed\u30b8\u30a7\u30af\u30bf\u30fc\u306b\u3088\u308b\u30b9\u30af\u30ea\u30fc\u30f3\u6295\u5f71",
    "AR/XR equipment": "AR/XR\u6a5f\u5668",
    "VR equipment": "VR\u6a5f\u5668",
    "Projection Mapping": "\u30d7\u30ed\u30b8\u30a7\u30af\u30b7\u30e7\u30f3\u30de\u30c3\u30d4\u30f3\u30b0",
    "Reciprocity and reality": "\u53cc\u65b9\u5411\u6027\u3068\u73fe\u5b9f\u611f",
    "Equipment related to virtual environment": "\u30d0\u30fc\u30c1\u30e3\u30eb\u74b0\u5883\u306b\u95a2\u3059\u308b\u6a5f\u5668",
    "Virtual environment output device type": "\u30d0\u30fc\u30c1\u30e3\u30eb\u74b0\u5883\u51fa\u529b\u30c7\u30d0\u30a4\u30b9\u306e\u7a2e\u985e",
    "Training content and types of virtual environment": "\u30c8\u30ec\u30fc\u30cb\u30f3\u30b0\u5185\u5bb9\u3068\u30d0\u30fc\u30c1\u30e3\u30eb\u74b0\u5883\u306e\u7a2e\u985e",
    "iClone": "iClone",
    " Motion LIVE": " Motion LIVE",
    "vSim": "vSim",
    "\u00ae for Nursing": "\u00ae for Nursing",
    "mixed reality": "\u8907\u5408\u73fe\u5b9f",
    "CG": "CG",
    "Metaverse-Based Clinical Training in Nursing Education: ": "\u770b\u8b77\u6559\u80b2\u306b\u304a\u3051\u308b\u30e1\u30bf\u30d0\u30fc\u30b9\u3092\u6d3b\u7528\u3057\u305f\u81e8\u5e8a\u30c8\u30ec\u30fc\u30cb\u30f3\u30b0\uff1a",
    "Utilization of Computer Graphics and Generative": "\u30b3\u30f3\u30d4\u30e5\u30fc\u30bf\u30b0\u30e9\u30d5\u30a3\u30c3\u30af\u30b9\u3068\u751f\u6210AI\u306e",

    # Slide 17 / 22
    "Maternal Nursing": "\u6bcd\u6027\u770b\u8b77",
    "Adult Nursing": "\u6210\u4eba\u770b\u8b77",
    "Home Nursing": "\u5728\u5b85\u770b\u8b77",
    "Incubator for premature babies": "\u672a\u719f\u5150\u7528\u4fdd\u80b2\u5668",
    "Home-visit nursing": "\u8a2a\u554f\u770b\u8b77",
    "sickroom": "\u75c5\u5ba4",

    # Slide 20 / 33 / 34 / 35 - Product name labels on slides
    "UNITY": "UNITY",
    "TWINMOTION": "TWINMOTION",

    # Slide 27
    "PowerPoint 3-screen video production": "PowerPoint 3\u753b\u9762\u6620\u50cf\u5236\u4f5c",
    "Since it is a projector, there are some issues with brightness and resolution.": "\u30d7\u30ed\u30b8\u30a7\u30af\u30bf\u30fc\u306e\u305f\u3081\u3001\u660e\u308b\u3055\u3068\u89e3\u50cf\u5ea6\u306b\u8ab2\u984c\u3042\u308a",
    "In large LCD display\nStart researching VR construction!": "\u5927\u578b\u6db2\u6676\u30c7\u30a3\u30b9\u30d7\u30ec\u30a4\u3067\nVR\u69cb\u7bc9\u306e\u7814\u7a76\u3092\u958b\u59cb\uff01",
    "Built a system that allows several people to experience VR at once with 3 prime projectors + 3 rear projection screens.": "3\u53f0\u306e\u30d7\u30ed\u30b8\u30a7\u30af\u30bf\u30fc\uff0b3\u9762\u30ea\u30a2\u30d7\u30ed\u30b8\u30a7\u30af\u30b7\u30e7\u30f3\u30b9\u30af\u30ea\u30fc\u30f3\u3067\u8907\u6570\u4eba\u304c\u540c\u6642\u306bVR\u4f53\u9a13\u3067\u304d\u308b\u30b7\u30b9\u30c6\u30e0\u3092\u69cb\u7bc9",
    "Three screens output from a three-screen spanned computer to three projectors": "3\u753b\u9762\u30b9\u30d1\u30f3PC\u304b\u30893\u53f0\u306e\u30d7\u30ed\u30b8\u30a7\u30af\u30bf\u30fc\u30783\u753b\u9762\u51fa\u529b",
    "Projection from 3 projectors to 3 rear projection large screens": "3\u53f0\u306e\u30d7\u30ed\u30b8\u30a7\u30af\u30bf\u30fc\u304b\u30893\u9762\u30ea\u30a2\u30d7\u30ed\u30b8\u30a7\u30af\u30b7\u30e7\u30f3\u5927\u578b\u30b9\u30af\u30ea\u30fc\u30f3\u3078\u6295\u5f71",
    "CAVE-type VR construction completed": "CAVE\u578bVR\u69cb\u7bc9\u5b8c\u4e86",
    "Consideration of large displays": "\u5927\u578b\u30c7\u30a3\u30b9\u30d7\u30ec\u30a4\u306e\u691c\u8a0e",
    "It turned out that it was still difficult to handle because it was heavy and took up a lot of space.": "\u91cd\u91cf\u304c\u3042\u308a\u5834\u6240\u3092\u53d6\u308b\u305f\u3081\u3001\u53d6\u308a\u6271\u3044\u304c\u56f0\u96e3\u3067\u3042\u308b\u3053\u3068\u304c\u5224\u660e",
    "Purchased 3 LG 75-inch 4K displays for about 400,000 yen (including stand)": "LG 75\u30a4\u30f3\u30c14K\u30c7\u30a3\u30b9\u30d7\u30ec\u30a43\u53f0\u3092\u7d0440\u4e07\u5186\uff08\u30b9\u30bf\u30f3\u30c9\u8fbc\u307f\uff09\u3067\u8cfc\u5165",
    "Brighter and clearer than a projector": "\u30d7\u30ed\u30b8\u30a7\u30af\u30bf\u30fc\u3088\u308a\u660e\u308b\u304f\u9bae\u660e",

    # Slide 31
    "Phased Learning Program": "\u6bb5\u968e\u7684\u5b66\u7fd2\u30d7\u30ed\u30b0\u30e9\u30e0",
    "Basic Level": "\u57fa\u790e\u30ec\u30d9\u30eb",
    "Observation & Communication": "\u89b3\u5bdf\u3068\u30b3\u30df\u30e5\u30cb\u30b1\u30fc\u30b7\u30e7\u30f3",
    "\u2022 Basic conversation": "\u2022 \u57fa\u672c\u7684\u306a\u4f1a\u8a71",
    "\u2022 Vital signs": "\u2022 \u30d0\u30a4\u30bf\u30eb\u30b5\u30a4\u30f3",
    "\u2022 Symptom observation": "\u2022 \u75c7\u72b6\u89b3\u5bdf",
    "Intermediate Level": "\u4e2d\u7d1a\u30ec\u30d9\u30eb",
    "Judgment & Practice": "\u5224\u65ad\u3068\u5b9f\u8df5",
    "\u2022 Care planning": "\u2022 \u30b1\u30a2\u8a08\u753b",
    "\u2022 Equipment operation": "\u2022 \u6a5f\u5668\u64cd\u4f5c",
    "\u2022 Emergency response": "\u2022 \u7dca\u6025\u5bfe\u5fdc",
    "Advanced Level": "\u4e0a\u7d1a\u30ec\u30d9\u30eb",
    "Integration & Application": "\u7d71\u5408\u3068\u5fdc\u7528",
    "\u2022 Multiple patients": "\u2022 \u8907\u6570\u60a3\u8005",
    "\u2022 Team collaboration": "\u2022 \u30c1\u30fc\u30e0\u9023\u643a",
    "\u2022 Ethical dilemmas": "\u2022 \u502b\u7406\u7684\u30b8\u30ec\u30f3\u30de",
    "AI-Powered Individualized Learning": "AI\u99c6\u52d5\u306e\u500b\u5225\u5316\u5b66\u7fd2",
    "Weakness identification | Scenario adjustment | Performance tracking": "\u5f31\u70b9\u306e\u7279\u5b9a | \u30b7\u30ca\u30ea\u30aa\u8abf\u6574 | \u30d1\u30d5\u30a9\u30fc\u30de\u30f3\u30b9\u8ffd\u8de1",

    # Slide 32
    "Implementation Challenges & Solutions": "\u5b9f\u88c5\u4e0a\u306e\u8ab2\u984c\u3068\u89e3\u6c7a\u7b56",
    "Technical Challenges": "\u6280\u8853\u7684\u8ab2\u984c",
    "\u2022 AI hallucination phenomena": "\u2022 AI\u30cf\u30eb\u30b7\u30cd\u30fc\u30b7\u30e7\u30f3\u73fe\u8c61",
    "\u2022 System integration complexity": "\u2022 \u30b7\u30b9\u30c6\u30e0\u7d71\u5408\u306e\u8907\u96d1\u3055",
    "\u2022 Data privacy protection": "\u2022 \u30c7\u30fc\u30bf\u30d7\u30e9\u30a4\u30d0\u30b7\u30fc\u4fdd\u8b77",
    "\u2022 Hardware requirements": "\u2022 \u30cf\u30fc\u30c9\u30a6\u30a7\u30a2\u8981\u4ef6",
    "Technical Solutions": "\u6280\u8853\u7684\u89e3\u6c7a\u7b56",
    "\u2022 Multi-layer verification systems": "\u2022 \u591a\u5c64\u691c\u8a3c\u30b7\u30b9\u30c6\u30e0",
    "\u2022 Phased implementation approach": "\u2022 \u6bb5\u968e\u7684\u5c0e\u5165\u30a2\u30d7\u30ed\u30fc\u30c1",
    "\u2022 Privacy protection frameworks": "\u2022 \u30d7\u30e9\u30a4\u30d0\u30b7\u30fc\u4fdd\u8b77\u30d5\u30ec\u30fc\u30e0\u30ef\u30fc\u30af",
    "\u2022 Cloud-based infrastructure": "\u2022 \u30af\u30e9\u30a6\u30c9\u30d9\u30fc\u30b9\u306e\u30a4\u30f3\u30d5\u30e9",
    "Organizational Challenges": "\u7d44\u7e54\u7684\u8ab2\u984c",
    "\u2022 Faculty training needs": "\u2022 \u6559\u54e1\u7814\u4fee\u306e\u5fc5\u8981\u6027",
    "\u2022 Initial investment costs": "\u2022 \u521d\u671f\u6295\u8cc7\u30b3\u30b9\u30c8",
    "\u2022 Change management": "\u2022 \u5909\u66f4\u7ba1\u7406",
    "Organizational Solutions": "\u7d44\u7e54\u7684\u89e3\u6c7a\u7b56",
    "\u2022 Phased training programs": "\u2022 \u6bb5\u968e\u7684\u7814\u4fee\u30d7\u30ed\u30b0\u30e9\u30e0",
    "\u2022 Strategic investment planning": "\u2022 \u6226\u7565\u7684\u6295\u8cc7\u8a08\u753b",
    "\u2022 Stakeholder engagement": "\u2022 \u30b9\u30c6\u30fc\u30af\u30db\u30eb\u30c0\u30fc\u306e\u95a2\u4e0e",

    # Slide 37
    "1. Far away from the palace, far away from the guidance of the introduction": "1. \u9060\u9694\u5b9f\u7fd2\u6307\u5c0e\u306e\u5c0e\u5165",
    "Hospital A and University A Relay": "A\u75c5\u9662\u3068A\u5927\u5b66\u306e\u4e2d\u7d99",
    "University A Hospital": "A\u5927\u5b66\u75c5\u9662",
    "Send and receive credit PC": "\u9001\u53d7\u4fe1\u7528PC",
    "LTE or ": "LTE\u307e\u305f\u306f",
    "WiFi": "WiFi",
    "Photography camera": "\u64ae\u5f71\u7528\u30ab\u30e1\u30e9",
    "A University": "A\u5927\u5b66",
    "Relay of clinical scenes": "\u81e8\u5e8a\u5834\u9762\u306e\u4e2d\u7d99",
    "Distance Instruction for Students": "\u5b66\u751f\u3078\u306e\u9060\u9694\u6307\u5c0e",
    "Teleconferencing": "\u30c6\u30ec\u30d3\u4f1a\u8b70",
    "2. Introduction of cutting-edge technology into medical education": "2. \u533b\u7642\u6559\u80b2\u3078\u306e\u5148\u7aef\u6280\u8853\u306e\u5c0e\u5165",
    "3. Utilization of XR for practical training guidance": "3. \u5b9f\u7fd2\u6307\u5c0e\u3078\u306eXR\u6d3b\u7528",
    "4. Application to medical education using existing state-of-the-art technology": "4. \u65e2\u5b58\u306e\u6700\u5148\u7aef\u6280\u8853\u3092\u7528\u3044\u305f\u533b\u7642\u6559\u80b2\u3078\u306e\u5fdc\u7528",
    "Leveraging the Metaverse": "\u30e1\u30bf\u30d0\u30fc\u30b9\u306e\u6d3b\u7528",
    "Utilization of CG": "CG\u306e\u6d3b\u7528",
    "Tablet devices": "\u30bf\u30d6\u30ec\u30c3\u30c8\u7aef\u672b",
    "Large screen": "\u5927\u753b\u9762",
    "projector": "\u30d7\u30ed\u30b8\u30a7\u30af\u30bf\u30fc",
    "Large-format display": "\u5927\u578b\u30c7\u30a3\u30b9\u30d7\u30ec\u30a4",
    "A: Within the University\nReceiving video as a group": "A\uff1a\u5b66\u5185\n\u30b0\u30eb\u30fc\u30d7\u3067\u306e\u6620\u50cf\u53d7\u4fe1",
    "Metaverse\nplatform": "\u30e1\u30bf\u30d0\u30fc\u30b9\n\u30d7\u30e9\u30c3\u30c8\u30d5\u30a9\u30fc\u30e0",
    "CG software": "CG\u30bd\u30d5\u30c8\u30a6\u30a7\u30a2",
    "High-performance personal computers": "\u9ad8\u6027\u80fd\u30d1\u30bd\u30b3\u30f3",
    "On-campus Wi-Fi environment": "\u5b66\u5185Wi-Fi\u74b0\u5883",
    "Joint development of digital teaching materials": "\u30c7\u30b8\u30bf\u30eb\u6559\u6750\u306e\u5171\u540c\u958b\u767a",
    "VR experience at the University of F": "F\u5927\u5b66\u3067\u306eVR\u4f53\u9a13",
    "Integrated at the University of C": "C\u5927\u5b66\u3067\u306e\u7d71\u5408",
    "Leveraging XR\nHolding joint conferences, etc.": "XR\u306e\u6d3b\u7528\n\u5408\u540c\u30ab\u30f3\u30d5\u30a1\u30ec\u30f3\u30b9\u306e\u958b\u50ac\u7b49",
    "Provision of practical training content to desired universities": "\u5e0c\u671b\u3059\u308b\u5927\u5b66\u3078\u306e\u5b9f\u7fd2\u30b3\u30f3\u30c6\u30f3\u30c4\u306e\u63d0\u4f9b",
    "VR Experience at E University": "E\u5927\u5b66\u3067\u306eVR\u4f53\u9a13",
    "MR experience at the University of D": "D\u5927\u5b66\u3067\u306eMR\u4f53\u9a13",
    "A: Within the University\nIndividual video reception": "A\uff1a\u5b66\u5185\n\u500b\u4eba\u3067\u306e\u6620\u50cf\u53d7\u4fe1",
    "VR devices ": "VR\u30c7\u30d0\u30a4\u30b9 ",
    "VR devices": "VR\u30c7\u30d0\u30a4\u30b9",
    "MR(": "MR(",
    "hololens2)": "HoloLens2)",
    "Utilization of HTML programs": "HTML\u30d7\u30ed\u30b0\u30e9\u30e0\u306e\u6d3b\u7528",
    "Leverage the capture function": "\u30ad\u30e3\u30d7\u30c1\u30e3\u6a5f\u80fd\u306e\u6d3b\u7528",
    "Use of 360-degree cameras": "360\u5ea6\u30ab\u30e1\u30e9\u306e\u6d3b\u7528",
    "CG production at University A": "A\u5927\u5b66\u3067\u306eCG\u5236\u4f5c",
    "XR production at University B": "B\u5927\u5b66\u3067\u306eXR\u5236\u4f5c",

    # Slide 38
    "Conclusions": "\u7d50\u8ad6",
    "Key Contributions": "\u4e3b\u8981\u306a\u8ca2\u732e",
    "\u2022 UE5 enables unprecedented clinical environment reproduction": "\u2022 UE5\u306b\u3088\u308b\u524d\u4f8b\u306e\u306a\u3044\u81e8\u5e8a\u74b0\u5883\u306e\u518d\u73fe",
    "\u2022 Generative AI creates dynamic, individualized learning experiences": "\u2022 \u751f\u6210AI\u306b\u3088\u308b\u30c0\u30a4\u30ca\u30df\u30c3\u30af\u3067\u500b\u5225\u5316\u3055\u308c\u305f\u5b66\u7fd2\u4f53\u9a13\u306e\u5275\u51fa",
    "\u2022 Efficient workflows make implementation feasible for institutions": "\u2022 \u52b9\u7387\u7684\u306a\u30ef\u30fc\u30af\u30d5\u30ed\u30fc\u306b\u3088\u308a\u6559\u80b2\u6a5f\u95a2\u3067\u306e\u5c0e\u5165\u304c\u5b9f\u73fe\u53ef\u80fd\u306b",
    "Future Impact": "\u5c06\u6765\u3078\u306e\u5f71\u97ff",
    "\u2022 Fundamental transformation of nursing education paradigm": "\u2022 \u770b\u8b77\u6559\u80b2\u30d1\u30e9\u30c0\u30a4\u30e0\u306e\u6839\u672c\u7684\u306a\u5909\u9769",
    "\u2022 Overcome constraints: safety, individualization, repeatability": "\u2022 \u5236\u7d04\u306e\u514b\u670d\uff1a\u5b89\u5168\u6027\u3001\u500b\u5225\u5316\u3001\u53cd\u5fa9\u53ef\u80fd\u6027",
    "\u2022 Contribute to qualitative improvement and widespread adoption": "\u2022 \u8cea\u7684\u5411\u4e0a\u3068\u666e\u53ca\u3078\u306e\u8ca2\u732e",
    "Nursing education metaverse: Not just technology, but educational transformation": "\u770b\u8b77\u6559\u80b2\u30e1\u30bf\u30d0\u30fc\u30b9\uff1a\u6280\u8853\u3060\u3051\u3067\u306a\u304f\u3001\u6559\u80b2\u306e\u5909\u9769",

    # Slide 39
    "Thank You": "\u3054\u6e05\u8074\u3042\u308a\u304c\u3068\u3046\u3054\u3056\u3044\u307e\u3057\u305f",
    "Questions & Discussion": "\u8cea\u7591\u5fdc\u7b54",
    "Contact Information": "\u9023\u7d61\u5148",
    "Takeshi Koike": "\u5c0f\u6c60 \u6b66\u55e3",
    "Seirei Christopher University": "\u8056\u96b7\u30af\u30ea\u30b9\u30c8\u30d5\u30a1\u30fc\u5927\u5b66",
    "Yukie Majima": "\u771f\u5d8b \u7531\u8cb4\u60e0",
    "Osaka Metropolitan University": "\u5927\u962a\u516c\u7acb\u5927\u5b66",
}

def translate_run(run):
    text = run.text
    if text in translations:
        run.text = translations[text]
        return True
    return False

count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if translate_run(run):
                        count += 1
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if translate_run(run):
                                count += 1

prs.save(out_path)
print(f"Translated {count} text runs")
print(f"Saved to: {out_path}")
